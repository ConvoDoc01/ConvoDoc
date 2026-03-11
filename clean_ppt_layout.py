import sys
from pptx import Presentation

input_file = sys.argv[1]
output_file = sys.argv[2]

prs = Presentation(input_file)

for slide in prs.slides:

    # Collect only text shapes
    text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]

    # Sort top to bottom
    text_shapes.sort(key=lambda s: s.top)

    grouped = []
    current_group = []

    for shape in text_shapes:

        if not shape.text.strip():
            continue

        if not current_group:
            current_group.append(shape)
            continue

        prev = current_group[-1]

        # Alignment tolerance
        same_left = abs(shape.left - prev.left) < 15000
        vertical_gap = abs(shape.top - (prev.top + prev.height))

        if same_left and vertical_gap < 30000:
            current_group.append(shape)
        else:
            grouped.append(current_group)
            current_group = [shape]

    if current_group:
        grouped.append(current_group)

    # Merge grouped shapes
    for group in grouped:

        if len(group) <= 1:
            continue

        first = group[0]
        combined_text = ""

        for shp in group:
            combined_text += shp.text.strip() + " "

        # Remove old shapes
        for shp in group:
            slide.shapes._spTree.remove(shp._element)

        # Create new merged textbox
        new_box = slide.shapes.add_textbox(
            first.left,
            first.top,
            first.width,
            sum(s.height for s in group)
        )

        new_box.text_frame.text = combined_text.strip()

prs.save(output_file)
